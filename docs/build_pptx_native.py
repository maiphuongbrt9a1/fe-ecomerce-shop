"""
Build a fully editable PPTX of the app-flow diagram using native PowerPoint
shapes. Each "node card" is grouped into a single selectable group so you can
drag the whole card around in PowerPoint without ungrouping.

Run:
    python build_pptx_native.py
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os

# ---------- canvas: 16:9 widescreen, mapped to SVG viewBox 1920×1080 ----------
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
VB_W = 1920
VB_H = 1080
EMU_PER_IN = 914400

SX = (SLIDE_W_IN * EMU_PER_IN) / VB_W
SY = (SLIDE_H_IN * EMU_PER_IN) / VB_H


def x(v): return Emu(int(v * SX))
def y(v): return Emu(int(v * SY))


# ---------- palette ----------
BG       = RGBColor(0xEF, 0xEA, 0xE2)
NODE_BG  = RGBColor(0xFF, 0xFF, 0xFF)
NODE_BD  = RGBColor(0xC9, 0xC2, 0xB5)
TITLE    = RGBColor(0x1A, 0x1A, 0x1A)
ROUTE    = RGBColor(0x6A, 0x6A, 0x6A)
ARROW    = RGBColor(0x2F, 0x2F, 0x2F)
ARROW_SOFT = RGBColor(0x6B, 0x6B, 0x6B)
LABEL_BG = RGBColor(0xFF, 0xFF, 0xFF)
LABEL_BD = RGBColor(0xC9, 0xC2, 0xB5)

STRIPE = {
    "store": RGBColor(0x2C, 0x5F, 0x5D),
    "user":  RGBColor(0x57, 0x4A, 0x75),
    "auth":  RGBColor(0xB0, 0x55, 0x38),
    "admin": RGBColor(0x7E, 0x31, 0x40),
    "staff": RGBColor(0x3F, 0x6B, 0x4E),
}

PANEL = {
    "store": RGBColor(0xD6, 0xCF, 0xB8),
    "user":  RGBColor(0xCF, 0xC8, 0xD6),
    "auth":  RGBColor(0xD8, 0xC8, 0xBC),
    "admin": RGBColor(0xD2, 0xC2, 0xC8),
    "staff": RGBColor(0xC8, 0xD4, 0xC8),
}


# ---------- helpers ----------
def add_rect(slide, X, Y, W, H, fill, border=None, line_w=0.75, dash=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X, Y, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(line_w)
        if dash:
            ln = shp.line._get_or_add_ln()
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'dash')
    shp.shadow.inherit = False
    return shp


def add_text(slide, X, Y, W, H, text, *, size=11, bold=False, italic=False,
             color=TITLE, align="left", font="Inter"):
    tb = slide.shapes.add_textbox(X, Y, W, H)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    if align == "center": p.alignment = PP_ALIGN.CENTER
    elif align == "right": p.alignment = PP_ALIGN.RIGHT
    run = p.add_run(); run.text = text
    f = run.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return tb


# ---------- grouping (no API in python-pptx — manipulate XML) ----------
_GID = [1000]  # monotonically increasing group id


def _next_id():
    _GID[0] += 1
    return _GID[0]


def group_shapes(slide, shapes, name="Group"):
    """
    Wrap given shape objects into a single <p:grpSp> on the slide.
    Children keep their absolute coordinates (chOff=off, chExt=ext).
    """
    if not shapes:
        return None
    spTree = slide.shapes._spTree

    # bounding box of children
    xs = [int(s.left) for s in shapes]
    ys = [int(s.top) for s in shapes]
    xe = [int(s.left) + int(s.width) for s in shapes]
    ye = [int(s.top) + int(s.height) for s in shapes]
    off_x, off_y = min(xs), min(ys)
    ext_cx, ext_cy = max(xe) - off_x, max(ye) - off_y

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    grp = etree.SubElement(spTree, f"{{{p_ns}}}grpSp")
    nvGrpSpPr = etree.SubElement(grp, f"{{{p_ns}}}nvGrpSpPr")
    cNvPr = etree.SubElement(nvGrpSpPr, f"{{{p_ns}}}cNvPr")
    cNvPr.set("id", str(_next_id()))
    cNvPr.set("name", name)
    etree.SubElement(nvGrpSpPr, f"{{{p_ns}}}cNvGrpSpPr")
    etree.SubElement(nvGrpSpPr, f"{{{p_ns}}}nvPr")

    grpSpPr = etree.SubElement(grp, f"{{{p_ns}}}grpSpPr")
    xfrm = etree.SubElement(grpSpPr, f"{{{a_ns}}}xfrm")
    off = etree.SubElement(xfrm, f"{{{a_ns}}}off")
    off.set("x", str(off_x)); off.set("y", str(off_y))
    ext = etree.SubElement(xfrm, f"{{{a_ns}}}ext")
    ext.set("cx", str(ext_cx)); ext.set("cy", str(ext_cy))
    chOff = etree.SubElement(xfrm, f"{{{a_ns}}}chOff")
    chOff.set("x", str(off_x)); chOff.set("y", str(off_y))
    chExt = etree.SubElement(xfrm, f"{{{a_ns}}}chExt")
    chExt.set("cx", str(ext_cx)); chExt.set("cy", str(ext_cy))

    # move child shapes into the group
    for s in shapes:
        el = s._element
        el.getparent().remove(el)
        grp.append(el)

    return grp


def add_node(slide, X, Y, W, H, title, route, stripe_key, dashed=False):
    """White card + stripe + title + route; all four grouped into one."""
    card = add_rect(slide, x(X), y(Y), x(W) - x(0), y(H) - y(0),
                    fill=NODE_BG, border=NODE_BD, line_w=0.75, dash=dashed)
    stripe = add_rect(slide, x(X), y(Y), x(6) - x(0), y(H) - y(0),
                      fill=STRIPE[stripe_key])
    title_tb = add_text(slide, x(X + 14), y(Y + 6),
                        x(W - 16) - x(0), y(22) - y(0),
                        title, size=11, bold=True, color=TITLE)
    route_tb = add_text(slide, x(X + 14), y(Y + H - 24),
                        x(W - 16) - x(0), y(20) - y(0),
                        route, size=9, color=ROUTE, font="Consolas")
    group_shapes(slide, [card, stripe, title_tb, route_tb], name=f"Node:{title}")


def add_connector(slide, x1, y1, x2, y2, *, soft=False, dash=False,
                  bent=False, mid_x=None, mid_y=None):
    color = ARROW_SOFT if soft else ARROW
    width = Pt(1.0 if soft else 1.25)

    def _style(c, with_arrow):
        line = c.line
        line.color.rgb = color
        line.width = width
        ln = line._get_or_add_ln()
        for tag in ("a:tailEnd", "a:headEnd"):
            for el in ln.findall(qn(tag)):
                ln.remove(el)
        if dash:
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'dash')
        if with_arrow:
            tail = etree.SubElement(ln, qn('a:tailEnd'))
            tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')

    if not bent:
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       x(x1), y(y1), x(x2), y(y2))
        _style(c, with_arrow=True)
        return [c]

    if mid_x is None and mid_y is None:
        mid_x = x2
    if mid_x is not None and mid_y is None:
        # horiz then vert then horiz to (x2,y2)
        s1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x(x1), y(y1), x(mid_x), y(y1))
        s2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x(mid_x), y(y1), x(mid_x), y(y2))
        s3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x(mid_x), y(y2), x(x2), y(y2))
        segs = [s1, s2, s3]
    else:
        # vert then horiz then vert to (x2,y2)
        s1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x(x1), y(y1), x(x1), y(mid_y))
        s2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x(x1), y(mid_y), x(x2), y(mid_y))
        s3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x(x2), y(mid_y), x(x2), y(y2))
        segs = [s1, s2, s3]

    for i, c in enumerate(segs):
        _style(c, with_arrow=(i == len(segs) - 1))
    return segs


def add_label(slide, X, Y, W, text, *, size=9):
    """Single textbox with white fill + thin border + italic black text."""
    H = 16
    tb = slide.shapes.add_textbox(x(X), y(Y), x(W) - x(0), y(H) - y(0))
    # fill + border on the textbox itself
    tb.fill.solid(); tb.fill.fore_color.rgb = LABEL_BG
    tb.line.color.rgb = LABEL_BD
    tb.line.width = Pt(0.5)
    tb.shadow.inherit = False

    tf = tb.text_frame
    tf.margin_left = Emu(int(0.04 * EMU_PER_IN))
    tf.margin_right = Emu(int(0.04 * EMU_PER_IN))
    tf.margin_top = Emu(int(0.005 * EMU_PER_IN))
    tf.margin_bottom = Emu(int(0.005 * EMU_PER_IN))
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    f = run.font
    f.name = "Inter"; f.size = Pt(size); f.italic = True; f.color.rgb = TITLE
    return tb


def add_panel(slide, X, Y, W, H, fill, label):
    panel = add_rect(slide, x(X), y(Y), x(W) - x(0), y(H) - y(0), fill=fill)
    title = add_text(slide, x(X + 16), y(Y + 8),
                     x(W - 32) - x(0), y(20) - y(0),
                     label, size=9, bold=True,
                     color=RGBColor(0x3A, 0x3A, 0x3A))
    group_shapes(slide, [panel, title], name=f"Panel:{label}")


# ---------- build ----------
prs = Presentation()
prs.slide_width  = Emu(int(SLIDE_W_IN * EMU_PER_IN))
prs.slide_height = Emu(int(SLIDE_H_IN * EMU_PER_IN))
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
add_rect(slide, x(0), y(0), x(VB_W), y(VB_H), fill=BG)

# title
add_text(slide, x(32), y(20), x(1500) - x(0), y(36) - y(0),
         "Sơ đồ luồng ứng dụng — e-commerce shop",
         size=20, bold=True, color=TITLE)
add_text(slide, x(32), y(58), x(1800) - x(0), y(20) - y(0),
         "Mỗi mũi tên là một hành động hoặc điều kiện chuyển trang · 26 trang ứng dụng + middleware phân quyền",
         size=10, color=RGBColor(0x55, 0x55, 0x55))

# panels
add_panel(slide,   16,  92, 1888, 218, PANEL["store"], "CỬA HÀNG & MUA HÀNG")
add_panel(slide,   16, 324, 1888, 118, PANEL["user"],  "TÀI KHOẢN NGƯỜI DÙNG")
add_panel(slide,   16, 456, 1888, 138, PANEL["auth"],  "XÁC THỰC")
add_panel(slide,   16, 608, 1296, 456, PANEL["admin"], "QUẢN TRỊ — ADMIN · /admin/*")
add_panel(slide, 1324, 608,  580, 456, PANEL["staff"], "NHÂN VIÊN — STAFF · /staff/*")

# ===== Storefront =====
storefront = [
    ( 40, 150, 200, 62, "Trang chủ",          "/homepage"),
    (290, 150, 200, 62, "Tìm kiếm",           "/search"),
    (540, 150, 200, 62, "Chi tiết sản phẩm",  "/product/[id]"),
    (790, 150, 200, 62, "Giỏ hàng",           "/cart"),
    (1040, 150, 200, 62, "Thanh toán",        "/checkout"),
    (1290, 150, 280, 62, "Chi tiết đơn hàng", "/profile/orders/[id]"),
    (1040, 238, 240, 62, "Trả về VNPay",      "/checkout/vnpay-return"),
]
for X, Y, W, H, t, r in storefront:
    add_node(slide, X, Y, W, H, t, r, "store")

# ===== User account =====
for X, Y, W, H, t, r in [
    ( 40, 360, 200, 62, "Hồ sơ",          "/profile"),
    (290, 360, 200, 62, "Thông báo",      "/profile/notifications"),
    (540, 360, 200, 62, "Danh sách đơn",  "/profile/orders"),
    (790, 360, 200, 62, "Voucher của tôi","/profile/vouchers"),
]:
    add_node(slide, X, Y, W, H, t, r, "user")

# ===== Auth =====
for X, Y, W, H, t, r in [
    (  40, 500, 200, 62, "Đăng nhập",        "/auth/login"),
    ( 290, 500, 200, 62, "Đăng ký",          "/auth/signup"),
    ( 540, 500, 200, 62, "Xác thực tài khoản","/auth/verify/[id]"),
    ( 790, 500, 200, 62, "Quên mật khẩu",    "/auth/forgot-password"),
    (1040, 500, 200, 62, "Đổi mật khẩu",     "/auth/change-password"),
]:
    add_node(slide, X, Y, W, H, t, r, "auth")

# ===== Admin =====
for X, Y, W, H, t, r in [
    (  40, 670, 200, 62, "Bảng điều khiển", "/admin"),
    ( 252, 670, 200, 62, "Quản lý đơn hàng","/admin/orders"),
    ( 464, 670, 200, 62, "Khách hàng",      "/admin/users"),
    ( 676, 670, 200, 62, "Danh mục",        "/admin/categories"),
    ( 888, 670, 200, 62, "Hộp thoại / Chat","/admin/chat"),
    (1100, 670, 200, 62, "Mã giảm giá",     "/admin/coupons"),

    (  40, 760, 200, 62, "DS sản phẩm",     "/admin/products/list"),
    ( 252, 760, 200, 62, "Thêm sản phẩm",   "/admin/products/add"),
    ( 464, 760, 200, 62, "Sửa sản phẩm",    "/admin/products/edit/[id]"),
    ( 676, 760, 200, 62, "Đánh giá SP",     "/admin/products/reviews"),
    ( 888, 760, 200, 62, "Thư viện ảnh SP", "/admin/products/media"),
    (1100, 760, 200, 62, "Thương hiệu",     "/admin/brands"),

    (  40, 850, 200, 62, "Phân quyền",      "/admin/authority"),
    ( 252, 850, 200, 62, "Vai trò",         "/admin/roles"),
    ( 464, 850, 200, 62, "Cửa hàng",        "/admin/shop"),
    ( 676, 850, 200, 62, "Giao dịch",       "/admin/transactions"),
    ( 888, 850, 200, 62, "Hồ sơ admin",     "/admin/profile"),
]:
    add_node(slide, X, Y, W, H, t, r, "admin")

# Sidebar admin (dashed border) — also grouped
add_node(slide, 1100, 850, 200, 62, "Sidebar Admin", "điều hướng nội bộ",
         "admin", dashed=True)

# ===== Staff =====
for X, Y, W, H, t, r in [
    (1340, 670, 180, 62, "Bảng điều khiển", "/staff"),
    (1530, 670, 180, 62, "Quản lý đơn hàng","/staff/orders"),
    (1720, 670, 180, 62, "Khách hàng",      "/staff/users"),

    (1340, 760, 180, 62, "Danh mục",        "/staff/categories"),
    (1530, 760, 180, 62, "Hộp thoại / Chat","/staff/chat"),
    (1720, 760, 180, 62, "Mã giảm giá",     "/staff/coupons"),

    (1340, 850, 180, 62, "DS sản phẩm",     "/staff/products/list"),
    (1530, 850, 180, 62, "Đánh giá SP",     "/staff/products/reviews"),
    (1720, 850, 180, 62, "Hồ sơ staff",     "/staff/profile"),
]:
    add_node(slide, X, Y, W, H, t, r, "staff")


# =================== ARROWS ===================

# Storefront horizontal
for (a, b) in [(240, 290), (490, 540), (740, 790), (990, 1040), (1240, 1290)]:
    add_connector(slide, a, 181, b, 181)
add_label(slide, 241, 167,  48, "tìm kiếm")
add_label(slide, 490, 167,  50, "chọn SP")
add_label(slide, 736, 167,  58, "thêm giỏ")
add_label(slide, 988, 167,  54, "thanh toán")
add_label(slide,1228, 167,  62, "đặt thành công")

# Mua ngay: Chi tiết SP → Thanh toán (cong qua đỉnh)
add_connector(slide, 640, 150, 1140, 150, bent=True, mid_y=110)
add_label(slide, 846, 96, 138, 'người dùng nhấn "Mua ngay"')

# Thanh toán → VNPay
add_connector(slide, 1140, 212, 1140, 238)
add_label(slide, 1150, 220, 70, "chọn VNPay")

# VNPay → Chi tiết đơn (success)
add_connector(slide, 1280, 269, 1340, 269)
add_connector(slide, 1340, 269, 1340, 212)
add_label(slide, 1284, 256, 116, "thanh toán thành công")

# VNPay → Giỏ hàng (failed)
add_connector(slide, 1040, 269, 900, 269)
add_connector(slide, 900, 269, 900, 212)
add_label(slide, 908, 256, 100, "thanh toán thất bại")

# Trang chủ → Hồ sơ
add_connector(slide, 140, 212, 140, 360)
add_label(slide, 148, 276, 186, 'nhấn avatar → "Thông tin cá nhân"')

# Hồ sơ → tabs
for (a, b) in [(240, 391), (490, 391), (740, 391)]:
    add_connector(slide, a, 391, a + 50, 391)
add_label(slide, 248, 376, 178, "các tab trong trang Hồ sơ")

# Đơn hàng → Chi tiết đơn (lên trên)
add_connector(slide, 640, 360, 1410, 212, bent=True, mid_y=332)
add_label(slide, 980, 318, 142, "nhấn vào một đơn hàng")

# Voucher → Tìm kiếm
add_connector(slide, 890, 360, 390, 212, bent=True, mid_y=320)
add_label(slide, 582, 306, 116, 'nhấn "Dùng ngay"')

# Thanh toán → Đăng nhập (dashed)
add_connector(slide, 1100, 212, 210, 500, bent=True, mid_y=448, dash=True)
add_label(slide, 586, 441, 172, "người dùng chưa đăng nhập")

# ----- AUTH arrows -----
add_connector(slide, 240, 531, 290, 531)
add_label(slide, 240, 516, 106, 'nhấn "Đăng ký"')

add_connector(slide, 490, 531, 540, 531)
add_label(slide, 486, 516, 120, "đăng ký thành công")

# Xác thực → Đăng nhập (loop dưới)
add_connector(slide, 640, 562, 140, 562, bent=True, mid_y=579)
add_label(slide, 332, 571, 116, "xác thực thành công")

# Đăng nhập → Quên MK
add_connector(slide, 260, 562, 885, 562, bent=True, mid_y=595)
add_label(slide, 498, 587, 170, 'nhấn "Quên mật khẩu"')

# Quên MK → Đổi MK
add_connector(slide, 990, 531, 1040, 531)
add_label(slide, 982, 516, 116, "nhập email hợp lệ")

# Đổi MK → Đăng nhập (long curve below)
add_connector(slide, 1140, 562, 182, 562, bent=True, mid_y=592)
add_label(slide, 826, 584, 120, "đổi mật khẩu xong")

# ----- After login: 3 branches -----
add_connector(slide, 40, 520, 40, 170, bent=True, mid_x=20)
add_label(slide, 10, 340, 120, "đăng nhập (USER)")

add_connector(slide, 120, 562, 140, 670, bent=True, mid_y=640)
add_label(slide, 128, 603, 160, "có vai trò ADMIN")

add_connector(slide, 240, 520, 1430, 670, bent=True, mid_y=520)
add_label(slide, 700, 506, 180, "có vai trò OPERATOR")

# ----- ADMIN internal -----
add_connector(slide, 240, 701, 252, 701, soft=True, dash=True)

add_connector(slide, 240, 791, 252, 791, soft=True)
add_label(slide, 236, 777, 64, 'nhấn "Thêm SP"')

add_connector(slide, 240, 810, 464, 791, soft=True, bent=True, mid_y=810)
add_label(slide, 290, 813, 86, 'nhấn "Sửa SP"')

add_connector(slide, 352, 760, 464, 760, soft=True, bent=True, mid_y=740)
add_label(slide, 360, 730, 142, "sau khi tạo SP thành công")

add_connector(slide, 776, 732, 140, 760, soft=True, bent=True, mid_y=750)
add_label(slide, 240, 737, 200, 'nhấn "Xem SP của danh mục"')

add_connector(slide, 664, 701, 888, 701, soft=True)
add_label(slide, 700, 687, 160, 'nhấn "Nhắn KH" → chat')

add_connector(slide, 1100, 881, 40, 732, soft=True, dash=True, bent=True, mid_y=920)
add_label(slide, 440, 912, 240, "sidebar dẫn tới mọi trang /admin/*")

# ----- STAFF internal -----
add_connector(slide, 1430, 732, 1430, 850, soft=True, dash=True)
add_label(slide, 1438, 775, 140, "điều hướng qua sidebar")


# ----- LEGEND -----
LX, LY = 32, 1045
add_text(slide, x(LX), y(LY), x(80) - x(0), y(16) - y(0),
         "Chú giải:", size=9, bold=True, color=RGBColor(0x55, 0x55, 0x55))

legend_items = [
    ( 92, "store", "Cửa hàng"),
    (186, "user",  "Tài khoản"),
    (282, "auth",  "Xác thực"),
    (372, "admin", "Admin"),
    (446, "staff", "Staff"),
]
for lx_off, key, label in legend_items:
    sw = add_rect(slide, x(LX + lx_off), y(LY + 2),
                  x(12) - x(0), y(12) - y(0), fill=STRIPE[key])
    txt = add_text(slide, x(LX + lx_off + 16), y(LY),
                   x(80) - x(0), y(16) - y(0),
                   label, size=9, color=RGBColor(0x55, 0x55, 0x55))
    group_shapes(slide, [sw, txt], name=f"Legend:{label}")


def legend_line(LX, LY, dash=False, soft=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        x(LX), y(LY + 6), x(LX + 38), y(LY + 6))
    c.line.color.rgb = ARROW_SOFT if soft else ARROW
    c.line.width = Pt(1.0 if soft else 1.25)
    if dash:
        ln = c.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')

legend_line(LX + 540, LY)
add_text(slide, x(LX + 586), y(LY), x(150) - x(0), y(16) - y(0),
         "điều hướng chính", size=9, color=RGBColor(0x55, 0x55, 0x55))

legend_line(LX + 720, LY, soft=True)
add_text(slide, x(LX + 766), y(LY), x(150) - x(0), y(16) - y(0),
         "điều hướng nội bộ", size=9, color=RGBColor(0x55, 0x55, 0x55))

legend_line(LX + 900, LY, dash=True)
add_text(slide, x(LX + 946), y(LY), x(200) - x(0), y(16) - y(0),
         "điều kiện / middleware", size=9, color=RGBColor(0x55, 0x55, 0x55))


OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "app-flow-diagram.pptx")
prs.save(OUT)
print(f"OK -> {OUT}")
