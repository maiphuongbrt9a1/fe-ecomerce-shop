from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image
import os

DOCS = os.path.dirname(os.path.abspath(__file__))
PNG  = os.path.join(DOCS, "app-flow-diagram.png")
PPTX = os.path.join(DOCS, "app-flow-diagram.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

img = Image.open(PNG)
img_w, img_h = img.size
slide_w, slide_h = prs.slide_width, prs.slide_height
slide_ratio = slide_w / slide_h
img_ratio = img_w / img_h

if img_ratio > slide_ratio:
    w = slide_w
    h = int(w / img_ratio)
else:
    h = slide_h
    w = int(h * img_ratio)

left = int((slide_w - w) / 2)
top  = int((slide_h - h) / 2)
slide.shapes.add_picture(PNG, left, top, width=w, height=h)

prs.save(PPTX)
print(f"OK -> {PPTX}")
